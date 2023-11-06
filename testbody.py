from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    
    full_text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            # Überprüfung, ob die Form eine Textbox, ein AutoShape oder eine Gruppe mit Text ist
            if hasattr(shape, "text"):
                full_text.append(shape.text_frame.text)
    
    return '\n'.join(full_text)

# Example usage:
pptx_path = '/Users/emil/Desktop/Project/Abschlusspräsentation Schalke hilft!.pptx'
text = extract_text_from_pptx(pptx_path)
print(text)



