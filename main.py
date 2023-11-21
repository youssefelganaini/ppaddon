from pptx import Presentation
from gpt import send_gpt

PPTX_PATH = "./Abschlusspräsentation Schalke hilft!.pptx"
OUTPUT_PATH = "./response.txt"

text_dict = {}  # dict that maps slide number to action title


def extract_text_from_pptx(pptx_path):
    # Load the presentation
    prs = Presentation(pptx_path)

    # Iterate through slides, then through slide elements to extract text
    count_slide = 0
    for slide in prs.slides:
        count_slide += 1
        # add another dict key
        text_dict[count_slide] = None

        # Get Action title
        if slide.shapes.title is not None:
            # add slide action title to dict key
            text_dict[count_slide] = slide.shapes.title.text

    return text_dict


# Add string for ChatGPT prompt
gpt_string = ""
text = extract_text_from_pptx(PPTX_PATH)
for elem in text:
    gpt_string += "Folie " + str(elem) + ": " + str(text[elem]) + "\n"


# save output of chatgpt return value
def save_to_txt_file(text, output_path):
    with open(output_path, "w") as file:
        file.write(text)


response = send_gpt(gpt_string)
save_to_txt_file(response["choices"][0]["message"]["content"], OUTPUT_PATH)
