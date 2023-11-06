from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    # Load the presentation
    prs = Presentation(pptx_path)
    
    # Iterate through slides, then through slide elements to extract text
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    
    return '\n'.join(full_text)

def save_to_txt_file(text, output_path):
    with open(output_path, 'w') as file:
        file.write(text)

# Example usage:
pptx_path = '/Users/emil/Desktop/Project/MainualReadDeck.pptx'
output_path = '/Users/emil/Desktop/Project/respnse.txt'

text = extract_text_from_pptx(pptx_path)
save_to_txt_file(text, output_path)